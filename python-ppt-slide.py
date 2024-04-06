from pptx import Presentation

from pptx.util import Inches

from transformers import pipeline

generator = pipeline("text-generation", model="openai-community/gpt2")


title = input("Enter title")

prompt = input("Describe the topic")

prompt2 = input("second page")

prompt3 = input("third page")

prompt4 = input("fourth page")

generated_text = generator(
    
    prompt,
    max_length=1000,
    num_return_sequences=1,          
    truncation=True,
)
generated_text = generated_text[0]['generated_text']

print(generated_text)

pr1 = Presentation()

slide1_register = pr1.slide_layouts[0]

slide1 = pr1.slides.add_slide(slide1_register)

title1 = slide1.shapes.title

subtitle = slide1.placeholders[1]

title1.text = title

subtitle.text = generated_text

slide2_register = pr1.slide_layouts[1]

slide2 = pr1.slides.add_slide(slide2_register)

slide3_register = pr1.slide_layouts[5]

slide3 = pr1.slides.add_slide(slide3_register)

slide4_register = pr1.slide_layouts[5]

slide4 = pr1.slides.add_slide(slide4_register)

title2 = slide2.shapes.title

title2.text = "Some bullet points"

bullet_point_box = slide2.shapes

bullet_points_lvl1 = bullet_point_box.placeholders[1]

bullet_points_lvl1.text = "PPT generator test"

bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()

bullet_points_lvl2.text = "to"

bullet_points_lvl2.level = 1

pr1.save("ppt-gen.pptx")