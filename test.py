from flask import Flask, request, render_template, send_file, redirect, url_for
from pptx import Presentation
from pptx.util import Inches
# from transformers import pipeline
from pptx.dml.color import RGBColor
import os
from pathlib import Path
import base64
import tempfile

app = Flask(__name__)

VIBGYOR_COLORS = {
    "violet": (148, 0, 211),
    "indigo": (75, 0, 130),
    "blue": (0, 0, 255),
    "green": (0, 255, 0),
    "yellow": (255, 255, 0),
    "orange": (255, 165, 0),
    "red": (255, 0, 0)
}

def generate_slides(title, slide_prompts, theme_color, existing_pptx_file):
    generator = pipeline("text-generation", model="openai-community/gpt2")

    try:
        if existing_pptx_file and os.path.exists(existing_pptx_file):
            pr = Presentation(existing_pptx_file)
        else:
            raise FileNotFoundError

        if theme_color:
            try:
                if theme_color.lower() in VIBGYOR_COLORS:
                    color = RGBColor(*VIBGYOR_COLORS[(theme_color.lower())])
                else:
                    color = RGBColor.from_string(theme_color)

                pr.slide_master.shapes[0].fill.solid()
                pr.slide_master.shapes[0].fill.fore_color.rgb = color
            except ValueError:
                print("Invalid theme color name.")

        generated_texts = []
        for prompt in slide_prompts:
            generated_text = generator(prompt, max_length=200, num_return_sequences=1, truncation=True)[0]['generated_text']
            generated_texts.append(generated_text)

        for i, text in enumerate(generated_texts):
            slide_layout = pr.slide_layouts[1]
            slide = pr.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"{title} - Page {i+1}"

            if len(slide.placeholders) >= 2:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = text
            else:
                print(f"No subtitle placeholder found on slide {i+1}. Text content not added.")

        edited_pptx_filename = os.path.splitext(existing_pptx_file)[0] + "_edited.pptx"
        pr.save(edited_pptx_filename)
        return edited_pptx_filename

    except FileNotFoundError:
        pr = Presentation()

        if theme_color:
            try:
                if theme_color.lower() in VIBGYOR_COLORS:
                    color = RGBColor(*VIBGYOR_COLORS[(theme_color.lower())])
                else:
                    color = RGBColor.from_string(theme_color)

                pr.slide_master.shapes[0].fill.solid()
                pr.slide_master.shapes[0].fill.fore_color.rgb = color
            except ValueError:
                print("Invalid theme color name.")

        generated_texts = []
        for prompt in slide_prompts:
            generated_text = generator(prompt, max_length=200, num_return_sequences=1, truncation=True)[0]['generated_text']
            generated_texts.append(generated_text)

        for i, text in enumerate(generated_texts):
            slide_layout = pr.slide_layouts[0]
            slide = pr.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"{title} - Page {i+1}"

            if len(slide.placeholders) >= 2:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = text
            else:
                print(f"No subtitle placeholder found on slide {i+1}. Text content not added.")

        generated_pptx_filename = "generated_presentation.pptx"
        pr.save(generated_pptx_filename)
        return generated_pptx_filename

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        title = request.form['title']
        num_slides = int(request.form['num_slides'])
        theme_color = request.form['theme_color']
        slide_prompts = [request.form[f'slide_prompt_{i}'] for i in range(1, num_slides+1)]
        uploaded_file = request.files.get('pptx_file')

        if uploaded_file:
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            uploaded_file.save(temp_file.name)
            pptx_path = generate_slides(title, slide_prompts, theme_color, temp_file.name)
        else:
            pptx_path = generate_slides(title, slide_prompts, theme_color, "")

        return send_file(pptx_path, as_attachment=True)

    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)
