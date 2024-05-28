import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from transformers import pipeline
from pptx.dml.color import RGBColor
import os
from pathlib import Path
import base64

VIBGYOR_COLORS = {
    "violet": (148, 0, 211),
    "indigo": (75, 0, 130),
    "blue": (0, 0, 255),
    "green": (0, 255, 0),
    "yellow": (255, 255, 0),
    "orange": (255, 165, 0),
    "red": (255, 0, 0)
}
BACKGROUND_IMAGES = [
    Path("BACKGROUND_IMAGES/5cc6cb46-9628-438d-8699-2bba18c8e5ba.png"),
    Path("BACKGROUND_IMAGES/67c4279a-99ba-4b18-8516-691cdd2c4554.png"),
    Path("BACKGROUND_IMAGES/3428144_60241.jpg"),
    Path("BACKGROUND_IMAGES/16264056_rm309-adj-03.jpg")
]

def generate_slides(title, slide_prompts, theme_color, existing_pptx_file):
    """Generates or edits a presentation based on user input.

    Args:
        title (str): Title of the presentation.
        slide_prompts (list): List of topics for each slide.
        theme_color (str): Desired theme color (VIBGYOR or RGB format).
        existing_pptx_file (str): Path to the existing PowerPoint template.
    """
    generator = pipeline("text-generation", model="openai-community/gpt2")

    try:
        if existing_pptx_file and os.path.exists(existing_pptx_file):
            pr = Presentation(existing_pptx_file)
            st.write("Editing existing PowerPoint template...")
        else:
            raise FileNotFoundError

        # Set theme color if provided
        if theme_color:
            try:
                if theme_color.lower() in VIBGYOR_COLORS:
                    color = RGBColor.from_triplet(*VIBGYOR_COLORS[(theme_color.lower())])
                else:
                    color = RGBColor.from_string(theme_color)

                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.solid()
                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.fore_color.rgb = color
            except ValueError:
                st.write("Invalid theme color name.")

        generated_texts = []
        for prompt in slide_prompts:
            generated_text = generator(prompt, max_length=200, num_return_sequences=1, truncation=True)[0]['generated_text']
            generated_texts.append(generated_text)

        # Add or edit slides and populate content
        for i, text in enumerate(generated_texts):
            slide_layout = pr.slide_master.slide_layouts[1]  # Adjust this index based on your template
            slide = pr.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"{title} - Page {i+1}"

            # Check if placeholder with index 1 exists on the current slide layout
            if len(slide.placeholders) >= 2:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = text
            else:
                st.warning(f"No subtitle placeholder found on slide {i+1}. Text content not added.")

        # Save the edited presentation
        edited_pptx_filename = os.path.splitext(existing_pptx_file)[0] + "_edited.pptx"
        pr.save(edited_pptx_filename)
        st.write("Presentation edited successfully!")
        st.write(f"Edited presentation saved as: {edited_pptx_filename}")

        # Return the path to the edited presentation
        return edited_pptx_filename

    except FileNotFoundError:
        st.write("Creating a new PowerPoint presentation...")
        pr = Presentation()

        if theme_color:
            try:
                if theme_color.lower() in VIBGYOR_COLORS:
                    color = RGBColor.from_triplet(*VIBGYOR_COLORS[(theme_color.lower())])
                else:
                    color = RGBColor.from_string(theme_color)

                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.solid()
                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.fore_color.rgb = color
            except ValueError:
                st.write("Invalid theme color name.")

        generated_texts = []
        for prompt in slide_prompts:
            generated_text = generator(prompt, max_length=200, num_return_sequences=1, truncation=True)[0]['generated_text']
            generated_texts.append(generated_text)

        for i, text in enumerate(generated_texts):
            slide_layout = pr.slide_master.slide_layouts[0]  # Adjust this index based on your template
            slide = pr.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"{title} - Page {i+1}"

            # Check if placeholder with index 1 exists on the current slide layout
            if len(slide.placeholders) >= 2:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = text
            else:
                st.warning(f"No subtitle placeholder found on slide {i+1}. Text content not added.")

        generated_pptx_filename = "generated_presentation.pptx"
        pr.save(generated_pptx_filename)
        st.write("Presentation generated successfully!")
        st.write(f"Generated presentation saved as: {generated_pptx_filename}")

        # Return the path to the generated presentation
        return generated_pptx_filename


    except FileNotFoundError:
        st.write("Creating a new PowerPoint presentation...")
        pr = Presentation()

        if theme_color:
            try:
                if theme_color.lower() in VIBGYOR_COLORS:
                    color = RGBColor.from_triplet(*VIBGYOR_COLORS[(theme_color.lower())])
                else:
                    color = RGBColor.from_string(theme_color)

                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.solid()
                pr.slide_master.designs[0].slide_layouts[0].shapes[0].fill.fore_color.rgb = color
            except ValueError:
                st.write("Invalid theme color name.")

        generated_texts = []
        for prompt in slide_prompts:
            generated_text = generator(prompt, max_length=200, num_return_sequences=1, truncation=True)[0]['generated_text']
            generated_texts.append(generated_text)

        for i, text in enumerate(generated_texts):
            slide = pr.slides.add_slide(pr.slide_layouts[0])
            title_shape = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]

            title_shape.text = f"{title} - Page {i+1}"
            subtitle_placeholder.text = text

        generated_pptx_filename = "generated_presentation.pptx"
        pr.save(generated_pptx_filename)
        st.write("Presentation generated successfully!")
        st.write(f"Generated presentation saved as: {generated_pptx_filename}")

        # Return the path to the generated presentation
        return generated_pptx_filename

def get_ppt_download_link(ppt_filename):
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

def main():
    st.title("PowerPoint Presentation Generator")
    title = st.text_input("Enter presentation title:")
    num_slides = st.number_input("How many slides do you want?", min_value=1, step=1)
    

    slide_prompts = []
    for i in range(num_slides):
        topic = st.text_input(f"Enter topic for slide {i+1}:")
        slide_prompts.append(topic)

    theme_color = st.text_input("Enter desired theme color (VIBGYOR or RGB format):")
    
    generate_button = st.button("Generate Presentation")
    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    if uploaded_file is not None:
        with open("temp.pptx", "wb") as f:
            f.write(uploaded_file.getvalue())

        pptx_path = generate_slides(title, slide_prompts, theme_color, "temp.pptx")
        if pptx_path:
            st.markdown(get_ppt_download_link(pptx_path), unsafe_allow_html=True)
        else:
            st.markdown(get_ppt_download_link("temp.pptx"), unsafe_allow_html=True)

if __name__ == "__main__":
    main()
