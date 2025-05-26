# app/services.py
import asyncio
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
import os

from . import schemas, pptx_utils, prompt_optimizer, llm_selector, llm_integrations, config

settings = config.settings

# --- Helper for Asynchronous Placeholder ---
async def generate_slide_headings(main_topic: str, num_slides: int) -> List[schemas.SlideHeading]:
    """
    Generates slide headings using an LLM.
    """
    prompt = prompt_optimizer.optimize_heading_generation_prompt(main_topic, num_slides)
    
    # Conceptual: Select and call LLM
    # For now, using a mock response
    # llm_api = llm_selector.select_llm_for_task("heading_generation")
    # raw_headings_text = await llm_api.generate_text(prompt, max_tokens=num_slides * 20) # Rough estimate
    
    print(f"DEBUG: Heading generation prompt: {prompt}") # For debugging

    # MOCK RESPONSE for headings
    await asyncio.sleep(0.1) # Simulate network delay
    raw_headings_text = "\n".join([f"{i+1}. Heading for '{main_topic}' - Part {i+1}" for i in range(num_slides)])
    
    # Parse headings (assuming LLM returns a numbered list)
    headings = []
    for i, line in enumerate(raw_headings_text.strip().split('\n')):
        if i >= num_slides: break
        try:
            heading_text = line.split('.', 1)[1].strip() if '.' in line else line.strip()
            headings.append(schemas.SlideHeading(id=i + 1, heading=heading_text))
        except IndexError:
            headings.append(schemas.SlideHeading(id=i + 1, heading=f"Default Heading {i+1}")) # Fallback
            
    # If fewer headings generated than requested, fill with defaults
    while len(headings) < num_slides:
        idx = len(headings) + 1
        headings.append(schemas.SlideHeading(id=idx, heading=f"Additional Topic {idx}"))
        
    return headings[:num_slides]


async def generate_content_for_slide(
    heading: str,
    main_topic: str,
    style_tone: str,
    content_format: str,
    max_tokens_per_slide: int
) -> str:
    """
    Generates content for a single slide.
    """
    prompt = prompt_optimizer.optimize_slide_content_prompt(
        heading, main_topic, style_tone, content_format, max_tokens_per_slide
    )
    print(f"DEBUG: Slide content prompt for '{heading}': {prompt}") # For debugging

    # Conceptual: Select and call LLM
    # For now, using a mock response
    # llm_api = llm_selector.select_llm_for_task("content_generation")
    # slide_content = await llm_api.generate_text(prompt, max_tokens=max_tokens_per_slide)
    
    # MOCK RESPONSE for slide content
    await asyncio.sleep(0.1) # Simulate network delay
    if content_format == "bullet_points":
        slide_content = f"- Key point 1 for {heading}\n- Key point 2 for {heading}\n- Another important detail about {heading}."
    elif content_format == "summary":
        slide_content = f"This is a brief summary regarding {heading}. It covers the essential aspects derived from the main topic of {main_topic}."
    else: # paragraph
        slide_content = f"This is a detailed paragraph explaining {heading}. It elaborates on the significance of this sub-topic within the broader context of {main_topic}. More information would be filled in here to reach the desired token count, discussing various facets and implications related to '{heading}'."
    
    # Simple token truncation for mock (real LLMs handle this)
    # actual_tokens = prompt_optimizer.count_tokens_simple(slide_content)
    # if actual_tokens > max_tokens_per_slide:
    #     words = slide_content.split()
    #     # This is a very rough way to truncate, real tokenizers are needed
    #     slide_content = " ".join(words[:max_tokens_per_slide // 2]) + "..." # Approximation
        
    return slide_content


async def generate_presentation_slides_async(
    job_id: str,
    details: schemas.FinalPresentationRequest,
    output_path: Path,
    job_store: Dict
):
    try:
        job_store[job_id] = {"status": "processing", "message": "Initializing presentation..."}
        
        template_path = pptx_utils.get_template_path(details.template_choice, details.uploaded_template_path)

        if template_path and template_path.exists():
            prs = Presentation(template_path)
            job_store[job_id]["message"] = f"Using template: {template_path.name}"
        else:
            prs = Presentation() # Create new if no valid template
            job_store[job_id]["message"] = "No valid template found, creating a new presentation."
            # (You might want to apply a default slide master or style here for new presentations)

        theme_rgb_color = pptx_utils.parse_theme_color(details.theme_color)
        if theme_rgb_color:
            pptx_utils.apply_theme_color_to_master(prs, theme_rgb_color)
            job_store[job_id]["message"] = "Applied theme color."
        
        # Identify a suitable slide layout for content (e.g., Title and Content)
        # This could be made more sophisticated, allowing user to choose or detecting based on template
        content_slide_layout_idx = 1 # Common index for "Title and Content"
        if not prs.slide_layouts or len(prs.slide_layouts) <= content_slide_layout_idx:
            content_slide_layout_idx = 0 # Fallback to first layout

        total_slides = len(details.final_headings)
        for i, slide_heading_text in enumerate(details.final_headings):
            job_store[job_id]["message"] = f"Generating content for slide {i+1}/{total_slides}: {slide_heading_text}"
            
            # Simulate fetching main_topic if needed, or pass it through session
            # For this example, assuming details.main_topic is available via session_id if needed
            # but FinalPresentationRequest doesn't have it. We'll assume it was implicitly passed
            # or we use a placeholder. For now, let's assume main_topic came from original session.
            # This implies SESSION_DATA must still be accessible or relevant parts passed to FinalPresentationRequest
            
            # To get main_topic, we'd need to adjust how session_id data is handled or passed.
            # For now, using a placeholder if not available.
            main_topic_for_slide = "the overall presentation topic" # Placeholder
            # A better way: Ensure main_topic is part of FinalPresentationRequest or retrieved from session

            slide_content_text = await generate_content_for_slide(
                slide_heading_text,
                main_topic_for_slide, # This needs to be correctly sourced
                details.style_tone,
                details.content_format,
                details.max_tokens_per_slide
            )
            
            # Placeholder mapping - can be made more sophisticated
            # For example, user could define this per template.
            # Default: title is shape.title or ph[0], content is ph[1]
            current_placeholder_map = details.placeholder_map.get(i) if details.placeholder_map else None


            pptx_utils.add_slide_with_content(
                prs,
                title_text=slide_heading_text,
                content_text=slide_content_text,
                slide_layout_idx=content_slide_layout_idx,
                placeholder_map=current_placeholder_map
            )
            await asyncio.sleep(0.05) # Small delay to allow other tasks

        prs.save(output_path)
        job_store[job_id] = {
            "status": "completed",
            "message": "Presentation generated successfully.",
            "filename": output_path.name,
            "download_url": f"/download/{output_path.name}" # Construct based on your routing
        }
        
        # Schedule cleanup for uploaded template if it was used and is temporary
        if details.template_choice == "upload" and details.uploaded_template_path:
            upload_path = Path(details.uploaded_template_path)
            # Ensure it's within a designated temporary upload area before deleting its parent
            if settings.GENERATED_PPTS_DIR / "uploads" in upload_path.parents:
                 asyncio.create_task(pptx_utils.cleanup_file(upload_path.parent, settings.GENERATED_PPT_TTL_SECONDS + 10)) # Delete the session upload folder

    except Exception as e:
        print(f"Error in generate_presentation_slides_async for job {job_id}: {e}")
        job_store[job_id] = {
            "status": "failed",
            "message": "An error occurred during presentation generation.",
            "error_detail": str(e)
        }
    finally:
        # Clean up main session data after job completion/failure
        if details.session_id in config.SESSION_DATA: # Assuming SESSION_DATA is accessible via config import
            del config.SESSION_DATA[details.session_id]