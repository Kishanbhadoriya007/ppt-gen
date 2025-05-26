# app/main.py
import secrets
import asyncio
import os
from pathlib import Path
from fastapi import FastAPI, Request, Form, File, UploadFile, HTTPException, BackgroundTasks
from fastapi.responses import HTMLResponse, FileResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from typing import List, Optional
import shutil
import uuid

from . import schemas, config, services, pptx_utils

settings = config.settings

app = FastAPI()

app.mount("/static", StaticFiles(directory=settings.STATIC_DIR), name="static")
templates = Jinja2Templates(directory=settings.TEMPLATES_DIR)

# In-memory store for job progress (replace with DB for production)
JOB_STORE = {} 
SESSION_DATA = {} # To store intermediate data like headings, template choice

@app.on_event("startup")
async def startup_event():
    if not settings.SERVER_TEMPLATES_DIR.exists():
        settings.SERVER_TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    if not settings.GENERATED_PPTS_DIR.exists():
        settings.GENERATED_PPTS_DIR.mkdir(parents=True, exist_ok=True)
    
    # Add dummy templates if none exist for testing
    if not list(settings.SERVER_TEMPLATES_DIR.glob("*.pptx")):
        try:
            # Create a simple dummy pptx file for testing
            from pptx import Presentation
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            title = slide.shapes.title
            title.text = "Dummy Template"
            prs.save(settings.SERVER_TEMPLATES_DIR / "dummy_template.pptx")
        except Exception as e:
            print(f"Could not create dummy template: {e}")


async def cleanup_file(file_path: Path, delay: int):
    await asyncio.sleep(delay)
    try:
        if file_path.exists():
            os.remove(file_path)
            print(f"Cleaned up {file_path}")
    except Exception as e:
        print(f"Error cleaning up file {file_path}: {e}")

@app.get("/", response_class=HTMLResponse)
async def get_index(request: Request):
    server_templates = pptx_utils.list_server_templates()
    return templates.TemplateResponse("index.html", {
        "request": request,
        "server_templates": server_templates,
        "default_slides": settings.DEFAULT_SLIDES,
        "max_slides": settings.MAX_SLIDES,
        "default_tokens": settings.DEFAULT_TOKENS_PER_SLIDE,
        "max_tokens": settings.MAX_TOKENS_PER_SLIDE
    })

@app.post("/generate-headings", response_class=HTMLResponse)
async def generate_headings_form(
    request: Request,
    main_topic: str = Form(...),
    num_slides: int = Form(settings.DEFAULT_SLIDES),
    template_choice: str = Form(...), # 'upload' or 'template_name.pptx'
    theme_color: Optional[str] = Form(None),
    max_tokens_per_slide: int = Form(settings.DEFAULT_TOKENS_PER_SLIDE),
    style_tone: str = Form("neutral"),
    content_format: str = Form("bullet_points"),
    pptx_template_file: Optional[UploadFile] = File(None)
):
    session_id = str(uuid.uuid4())
    uploaded_template_path = None

    if template_choice == "upload":
        if not pptx_template_file or pptx_template_file.filename == "":
            raise HTTPException(status_code=400, detail="PPTX template file must be uploaded if 'upload' is chosen.")
        if not pptx_template_file.filename.endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Invalid file type. Only .pptx files are allowed.")
        
        temp_upload_dir = settings.GENERATED_PPTS_DIR / "uploads" / session_id
        temp_upload_dir.mkdir(parents=True, exist_ok=True)
        uploaded_template_path = temp_upload_dir / pptx_template_file.filename
        
        try:
            with open(uploaded_template_path, "wb") as buffer:
                shutil.copyfileobj(pptx_template_file.file, buffer)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to save uploaded template: {str(e)}")
        finally:
            pptx_template_file.file.close()
    elif not (settings.SERVER_TEMPLATES_DIR / template_choice).exists():
        raise HTTPException(status_code=400, detail=f"Selected server template '{template_choice}' not found.")

    # Store session data
    SESSION_DATA[session_id] = {
        "main_topic": main_topic,
        "num_slides": num_slides,
        "template_choice": template_choice,
        "uploaded_template_path": str(uploaded_template_path) if uploaded_template_path else None,
        "server_template_name": template_choice if template_choice != "upload" else None,
        "theme_color": theme_color,
        "max_tokens_per_slide": max_tokens_per_slide,
        "style_tone": style_tone,
        "content_format": content_format
    }

    try:
        headings = await services.generate_slide_headings(main_topic, num_slides)
        SESSION_DATA[session_id]["initial_headings"] = headings # Store initially generated headings
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate headings: {str(e)}")

    return templates.TemplateResponse("configure_slides.html", {
        "request": request,
        "session_id": session_id,
        "main_topic": main_topic,
        "num_slides": num_slides,
        "headings": headings,
        "template_choice": template_choice,
        "theme_color": theme_color,
        "max_tokens_per_slide": max_tokens_per_slide,
        "style_tone": style_tone,
        "content_format": content_format
    })


@app.post("/create-presentation", response_class=HTMLResponse)
async def create_presentation_form(
    request: Request,
    background_tasks: BackgroundTasks,
    session_id: str = Form(...),
    final_headings: List[str] = Form(...) # Comes as a list of strings
):
    if session_id not in SESSION_DATA:
        raise HTTPException(status_code=404, detail="Session not found or expired.")

    session_info = SESSION_DATA[session_id]
    
    presentation_details = schemas.FinalPresentationRequest(
        session_id=session_id,
        final_headings=final_headings,
        template_choice=session_info["template_choice"],
        uploaded_template_path=session_info.get("uploaded_template_path"),
        server_template_name=session_info.get("server_template_name"),
        theme_color=session_info.get("theme_color"),
        max_tokens_per_slide=session_info["max_tokens_per_slide"],
        style_tone=session_info["style_tone"],
        content_format=session_info["content_format"]
    )

    job_id = f"ppt_job_{secrets.token_hex(8)}"
    JOB_STORE[job_id] = {"status": "processing", "message": "Generating presentation..."}

    output_filename = f"{session_info['main_topic'].replace(' ', '_')}_{job_id}.pptx"
    output_path = settings.GENERATED_PPTS_DIR / output_filename

    # Run generation in background
    background_tasks.add_task(
        services.generate_presentation_slides_async,
        job_id,
        presentation_details,
        output_path,
        JOB_STORE
    )
    
    # Clean up session data after initiating job
    # del SESSION_DATA[session_id] # Or use a TTL cache

    return templates.TemplateResponse("download.html", {
        "request": request,
        "job_id": job_id,
        "filename": output_filename,
        "status_url": app.url_path_for("get_job_status", job_id=job_id)
    })

@app.get("/status/{job_id}")
async def get_job_status(job_id: str):
    job = JOB_STORE.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    
    status_data = {"job_id": job_id, "status": job["status"], "message": job.get("message")}
    if job["status"] == "completed":
        filename = job.get("filename")
        if filename:
             status_data["download_url"] = app.url_path_for("download_file", filename=filename)
    elif job["status"] == "failed":
        status_data["error"] = job.get("error_detail", "Unknown error")
        
    return schemas.JobStatus(**status_data)


@app.get("/download/{filename}")
async def download_file(filename: str, background_tasks: BackgroundTasks):
    file_path = settings.GENERATED_PPTS_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found or expired.")
    
    # Schedule cleanup after download is initiated
    background_tasks.add_task(cleanup_file, file_path, settings.GENERATED_PPT_TTL_SECONDS)
    
    return FileResponse(
        path=file_path,
        filename=filename,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

# Placeholder for __init__.py
# app/__init__.py
# (This file can be empty)

# Placeholder for crud.py (for future MongoDB integration)
# app/crud.py
# from pymongo import MongoClient
# from .config import settings

# async def get_user_preferences(user_id: str):
#     # client = MongoClient(settings.MONGO_DATABASE_URL)
#     # db = client[settings.MONGO_DATABASE_NAME]
#     # user_prefs = db.user_preferences.find_one({"user_id": user_id})
#     # client.close()
#     # return user_prefs
#     pass

# async def save_generated_ppt_info(user_id: str, ppt_filename: str, topic: str):
#     pass