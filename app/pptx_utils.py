# app/pptx_utils.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path
from typing import List, Optional, Tuple, Dict

from .config import settings

VIBGYOR_COLORS = {
    "violet": (148, 0, 211), "indigo": (75, 0, 130), "blue": (0, 0, 255),
    "green": (0, 255, 0), "yellow": (255, 255, 0), "orange": (255, 165, 0),
    "red": (255, 0, 0)
}

def parse_theme_color(color_str: Optional[str]) -> Optional[RGBColor]:
    if not color_str:
        return None
    color_str_lower = color_str.lower()
    if color_str_lower in VIBGYOR_COLORS:
        return RGBColor(*VIBGYOR_COLORS[color_str_lower])
    try:
        if color_str_lower.startswith("#"):
             return RGBColor.from_string(color_str_lower[1:])
        return RGBColor.from_string(color_str_lower) # Handles hex like "RRGGBB"
    except ValueError:
        print(f"Warning: Invalid theme color '{color_str}'.")
        return None

def apply_theme_color_to_master(prs: Presentation, color: Optional[RGBColor]):
    if not color:
        return
    try:
        # This is a simplistic approach; real theme application is more complex.
        # Attempting to change the background of the first master's first layout.
        master = prs.slide_masters[0]
        
        # Try to change accent1 color (often used for titles or main elements)
        # This requires a deeper understanding of OOXML and python-pptx theme manipulation.
        # For a simple background color change of a common shape:
        if master.slide_layouts:
            first_layout = master.slide_layouts[0]
            if first_layout.shapes and hasattr(first_layout.shapes[0], 'fill'):
                 # Check if it's a background shape or a common large shape
                if first_layout.shapes[0].name.lower().startswith("background") or \
                   (first_layout.shapes[0].width > prs.slide_width * 0.8 and first_layout.shapes[0].height > prs.slide_height * 0.8) :
                    try:
                        fill = first_layout.shapes[0].fill
                        fill.solid()
                        fill.fore_color.rgb = color
                        print(f"Applied theme color to master's first layout shape.")
                    except Exception as e:
                        print(f"Could not apply color to specific shape: {e}")

        # A more robust way would be to modify theme.xml, which python-pptx has limited direct support for.
        # For now, we'll print a message if a more complex theme change is needed.
        print("Note: Advanced theme color changes might require manual template modification or deeper XML manipulation.")

    except Exception as e:
        print(f"Error applying theme color to master: {e}")


def add_slide_with_content(
    prs: Presentation,
    title_text: str,
    content_text: str,
    slide_layout_idx: int = 1, # Typically "Title and Content"
    placeholder_map: Optional[Dict[str, int]] = None # {"title": 0, "content": 1}
):
    try:
        slide_layout = prs.slide_layouts[slide_layout_idx]
    except IndexError:
        print(f"Warning: Slide layout index {slide_layout_idx} out of range. Using layout 0.")
        slide_layout = prs.slide_layouts[0] # Default to title slide if chosen layout is invalid
        
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = None
    content_placeholder = None

    if placeholder_map:
        title_ph_idx = placeholder_map.get("title")
        content_ph_idx = placeholder_map.get("content")
        try:
            if title_ph_idx is not None and len(slide.placeholders) > title_ph_idx:
                title_placeholder = slide.placeholders[title_ph_idx]
            if content_ph_idx is not None and len(slide.placeholders) > content_ph_idx:
                content_placeholder = slide.placeholders[content_ph_idx]
        except IndexError:
            print("Warning: Provided placeholder index in map is out of bounds.")
    
    # Fallback to default title and first body placeholder if map fails or not provided
    if not title_placeholder:
        if slide.shapes.title:
            title_placeholder = slide.shapes.title
        elif len(slide.placeholders) > 0: # Try first placeholder as title
            title_placeholder = slide.placeholders[0]
            
    if not content_placeholder: # Try second placeholder as content, or first if title took shapes.title
        if len(slide.placeholders) > 1:
             content_placeholder = slide.placeholders[1]
        elif len(slide.placeholders) > 0 and slide.shapes.title and slide.placeholders[0] != slide.shapes.title:
            # If title used slide.shapes.title, and placeholder[0] is different, it could be content
            content_placeholder = slide.placeholders[0]


    if title_placeholder:
        title_placeholder.text = title_text
    else:
        print(f"Warning: No title placeholder found or assignable on slide for '{title_text}'.")

    if content_placeholder:
        tf = content_placeholder.text_frame
        tf.clear() # Clear existing text
        p = tf.add_paragraph()
        p.text = content_text
        # Potentially apply formatting based on content_format (e.g. bullet points)
        # For simplicity, this example just adds the text.
        # To make actual bullet points, you'd iterate lines and set paragraph levels.
        # Example for bullet points (if content_text is pre-formatted with newlines for bullets):
        # tf.clear()
        # lines = content_text.split('\n')
        # for i, line in enumerate(lines):
        #     p = tf.add_paragraph()
        #     p.text = line.strip()
        #     if i > 0 or len(lines) > 1: # Simple bullet for subsequent lines or if it's the only line meant as a bullet
        #         p.level = 0 # 0 for first level bullet
        # tf.auto_size = True # Or MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        print(f"Warning: No content placeholder found or assignable on slide for '{title_text}'. Adding text box.")
        # Fallback: add a new text box if no suitable placeholder
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8.5)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = content_text

    return slide


def list_server_templates() -> List[str]:
    if settings.SERVER_TEMPLATES_DIR.exists():
        return sorted([f.name for f in settings.SERVER_TEMPLATES_DIR.glob("*.pptx")])
    return []

def get_template_path(template_choice: str, uploaded_template_path: Optional[str]) -> Optional[Path]:
    if template_choice == "upload":
        if uploaded_template_path and Path(uploaded_template_path).exists():
            return Path(uploaded_template_path)
        return None # Should have been caught earlier if file doesn't exist
    else: # Server template
        path = settings.SERVER_TEMPLATES_DIR / template_choice
        return path if path.exists() else None
        
def count_text_placeholders(prs: Presentation, layout_idx: int = 1) -> Dict[str, int]:
    """Counts title and body/content placeholders for a given layout."""
    counts = {"title": 0, "content_body": 0, "other": 0, "total_placeholders":0}
    if layout_idx >= len(prs.slide_layouts):
        return counts
        
    layout = prs.slide_layouts[layout_idx]
    counts["total_placeholders"] = len(layout.placeholders)
    
    for ph in layout.placeholders:
        if ph.is_placeholder: # Redundant check, but good practice
            if ph.name.lower().startswith("title") or ph.placeholder_format.type == MSO_SHAPE.TITLE:
                counts["title"] +=1
            # Placeholder type for content can vary; common ones are BODY, CONTENT, OBJECT
            elif ph.name.lower().startswith("content") or \
                 ph.name.lower().startswith("text") or \
                 ph.name.lower().startswith("body") or \
                 ph.placeholder_format.type in [MSO_SHAPE.BODY, MSO_SHAPE.CONTENT, MSO_SHAPE.OBJECT, MSO_SHAPE.TEXT_BOX]:
                counts["content_body"] +=1
            else:
                counts["other"] +=1
    return counts